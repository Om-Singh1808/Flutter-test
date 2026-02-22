"""
document_processor.py
Extracts text from PDF, DOCX, or PPTX files and splits into overlapping chunks.
Dispatches to the correct extractor based on file extension.
"""

import re
from typing import List

# ── Constants ────────────────────────────────────────────────────────────────
_CHUNK_SIZE = 500
_CHUNK_OVERLAP = 50
_SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx"}


# ── Text Extraction ───────────────────────────────────────────────────────────

def _extract_pdf(file_path: str) -> str:
    """Extract text from a PDF file using pypdf."""
    import pypdf

    reader = pypdf.PdfReader(file_path)
    pages: List[str] = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text)
    return "\n".join(pages)


def _extract_docx(file_path: str) -> str:
    """Extract text from a Word (.docx) file using python-docx."""
    import docx

    doc = docx.Document(file_path)
    paragraphs: List[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)
    return "\n".join(paragraphs)


def _extract_pptx(file_path: str) -> str:
    """Extract text from a PowerPoint (.pptx) file using python-pptx."""
    from pptx import Presentation

    prs = Presentation(file_path)
    slides_text: List[str] = []
    for slide in prs.slides:
        slide_parts: List[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_parts.append(text)
        if slide_parts:
            slides_text.append(" ".join(slide_parts))
    return "\n".join(slides_text)


# ── Chunker ───────────────────────────────────────────────────────────────────

def _split_into_chunks(
    text: str,
    chunk_size: int = _CHUNK_SIZE,
    overlap: int = _CHUNK_OVERLAP,
) -> List[str]:
    """Sliding-window splitter.  step = chunk_size - overlap."""
    if not text:
        return []

    step = chunk_size - overlap
    chunks: List[str] = []
    start = 0

    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start += step

    return chunks


# ── Public API ────────────────────────────────────────────────────────────────

def process_document(file_path: str, file_ext: str) -> List[str]:
    """
    Public entry-point for the document processing pipeline.

    Extracts text from the given file (PDF / DOCX / PPTX) and returns a list
    of overlapping text chunks ready for embedding.

    Args:
        file_path: Absolute path to the temporary uploaded file.
        file_ext:  Lowercased file extension including dot, e.g. '.pdf'.

    Returns:
        List of non-empty string chunks.

    Raises:
        ValueError: If the extension is not supported.
    """
    ext = file_ext.lower()
    if ext not in _SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"Unsupported file type '{ext}'. "
            f"Supported types: {', '.join(sorted(_SUPPORTED_EXTENSIONS))}"
        )

    print(
        f"[document_processor] Extracting text from: {file_path}  (type={ext})")

    if ext == ".pdf":
        raw = _extract_pdf(file_path)
    elif ext == ".docx":
        raw = _extract_docx(file_path)
    elif ext == ".pptx":
        raw = _extract_pptx(file_path)

    # Normalise whitespace
    clean = re.sub(r"\s+", " ", raw).strip()

    if not clean:
        print("[document_processor] WARNING: No text could be extracted.")
        return []

    print(f"[document_processor] Extracted {len(clean)} characters of text.")
    chunks = _split_into_chunks(clean)
    print(
        f"[document_processor] Created {len(chunks)} chunks "
        f"(size={_CHUNK_SIZE}, overlap={_CHUNK_OVERLAP})."
    )
    return chunks
